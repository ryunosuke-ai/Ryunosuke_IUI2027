"""Generate fully editable PowerPoint versions of BASiS Figures 1 and 2.

Every visible component is a native PowerPoint shape, text box, or connector.
The resulting slides do not contain flattened screenshots of the figures.
"""

from pathlib import Path

# Importing the compiled submodule first avoids a Python 3.13 lazy-import issue
# observed with the universal macOS lxml wheel.
from lxml import etree as _lxml_etree  # noqa: F401
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
OUTPUT = HERE / "BASiS_figures_editable.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

FONT = "Aptos"
FONT_DISPLAY = "Aptos Display"
FONT_MATH = "Cambria Math"

INK = "20313F"
MUTED = "667784"
LINE = "9AAAB5"
BLUE = "2A6F97"
BLUE_FILL = "EAF3F8"
BLUE_GROUP = "F4F9FC"
TEAL = "2A9D8F"
TEAL_FILL = "E8F6F3"
ORANGE = "D97745"
ORANGE_FILL = "FFF1E8"
PANEL_BG = "F8FAFB"
WHITE = "FFFFFF"


def rgb(hex_color):
    return RGBColor.from_string(hex_color)


def set_shape_name(shape, name):
    try:
        shape.name = name
    except (AttributeError, ValueError):
        pass
    return shape


def set_line_end(line_format, arrow=True):
    """Add an editable DrawingML arrowhead to the end of a connector."""
    if not arrow:
        return
    line_xml = line_format._get_or_add_ln()
    for child in list(line_xml):
        if child.tag.endswith("tailEnd"):
            line_xml.remove(child)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line_xml.append(tail)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=12,
    color=INK,
    bold=False,
    font=FONT,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin=0.02,
    name=None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        set_shape_name(shape, name)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    lines = text.split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = 0.92
        for run in paragraph.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)
    return shape


def add_box(
    slide,
    x,
    y,
    w,
    h,
    fill=WHITE,
    border=LINE,
    border_width=1.25,
    radius=True,
    name=None,
):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if radius
        else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if name:
        set_shape_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(border)
    shape.line.width = Pt(border_width)
    return shape


def add_panel(slide, x, y, w, h, step, title, badge_color, name):
    add_box(
        slide,
        x,
        y,
        w,
        h,
        fill=PANEL_BG,
        border=LINE,
        border_width=1.2,
        name=f"{name} panel",
    )
    add_text(
        slide,
        x + 0.22,
        y + 0.16,
        w - 0.44,
        0.30,
        title,
        size=9.4,
        bold=True,
        font=FONT_DISPLAY,
        name=f"{name} title",
    )
    add_badge(slide, x - 0.06, y - 0.06, step, badge_color, f"{name} step")


def add_badge(slide, x, y, number, color, name):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.36), Inches(0.36)
    )
    set_shape_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = Inches(0)
    frame.margin_top = frame.margin_bottom = Inches(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = str(number)
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = rgb(WHITE)
    return shape


def add_connector(
    slide,
    x1,
    y1,
    x2,
    y2,
    color=INK,
    width=1.5,
    arrow=True,
    dash=None,
    name=None,
):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    if name:
        set_shape_name(connector, name)
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if dash:
        connector.line.dash_style = dash
    set_line_end(connector.line, arrow)
    return connector


def add_polyline(slide, points, color=INK, width=1.5, arrow=True, name="route"):
    for index in range(len(points) - 1):
        add_connector(
            slide,
            *points[index],
            *points[index + 1],
            color=color,
            width=width,
            arrow=arrow and index == len(points) - 2,
            name=f"{name} {index + 1}",
        )


def add_chip(
    slide,
    x,
    y,
    w,
    h,
    text,
    color,
    fill,
    size=9,
    name=None,
):
    add_box(
        slide,
        x,
        y,
        w,
        h,
        fill=fill,
        border=color,
        border_width=1.05,
        name=f"{name} box" if name else None,
    )
    return add_text(
        slide,
        x + 0.04,
        y + 0.02,
        w - 0.08,
        h - 0.04,
        text,
        size=size,
        color=color,
        font=FONT_MATH if any(c in text for c in "ₜ₋₁∑̂") else FONT,
        name=f"{name} text" if name else None,
    )


def add_line_label(
    slide,
    x,
    y,
    w,
    h,
    text,
    color=MUTED,
    size=8.3,
    name=None,
):
    add_box(
        slide,
        x,
        y,
        w,
        h,
        fill=WHITE,
        border=WHITE,
        border_width=0.1,
        radius=True,
        name=f"{name} background" if name else None,
    )
    return add_text(
        slide,
        x,
        y,
        w,
        h,
        text,
        size=size,
        color=color,
        name=name,
    )


def add_bar(slide, x, y, width, color, label, name):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(0.16),
    )
    set_shape_name(shape, f"{name} bar")
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)
    add_text(
        slide,
        x + width + 0.08,
        y - 0.05,
        0.42,
        0.25,
        label,
        size=8.2,
        align=PP_ALIGN.LEFT,
        name=f"{name} value",
    )


def add_database_icon(slide, x, y, color, name):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CAN, Inches(x), Inches(y), Inches(0.34), Inches(0.30)
    )
    set_shape_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.color.rgb = rgb(color)


def add_layers_icon(slide, x, y, color, name):
    for index, offset in enumerate((0.00, 0.09, 0.18)):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.DIAMOND,
            Inches(x),
            Inches(y + offset),
            Inches(0.34),
            Inches(0.20),
        )
        set_shape_name(shape, f"{name} {index + 1}")
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(color)
        shape.line.color.rgb = rgb(WHITE)
        shape.line.width = Pt(0.4)


def set_slide_background(slide, color=WHITE):
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = rgb(color)


def build_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)

    # Phase containers
    add_box(
        slide,
        0.24,
        0.22,
        7.55,
        3.18,
        fill=BLUE_GROUP,
        border="B9D3E2",
        border_width=1.2,
        name="Phase A background",
    )
    add_box(
        slide,
        0.24,
        3.58,
        12.84,
        3.46,
        fill=PANEL_BG,
        border="B7DDD8",
        border_width=1.2,
        name="Phase B background",
    )
    add_chip(
        slide,
        0.50,
        0.41,
        2.64,
        0.34,
        "A   STYLE MODEL INDUCTION",
        BLUE,
        BLUE_FILL,
        size=9,
        name="Phase A label",
    )
    add_chip(
        slide,
        0.50,
        3.76,
        3.11,
        0.34,
        "B   DATA SELECTION & ADAPTATION",
        TEAL,
        TEAL_FILL,
        size=9,
        name="Phase B label",
    )

    # Corpus and module cards
    add_box(slide, 0.55, 0.96, 1.62, 2.05, WHITE, LINE, 1.35, name="Small corpus")
    add_database_icon(slide, 1.19, 1.14, BLUE, "Small corpus database icon")
    add_text(
        slide,
        0.72,
        1.53,
        1.28,
        0.25,
        "STYLE REFERENCE",
        size=8.1,
        color=INK,
        name="Small corpus role",
    )
    add_text(
        slide,
        0.69,
        1.82,
        1.34,
        0.93,
        "Small,\nhigh-quality\ndialogue corpus",
        size=14.2,
        bold=True,
        font=FONT_DISPLAY,
        name="Small corpus label",
    )

    add_box(
        slide,
        2.52,
        1.16,
        2.16,
        1.58,
        BLUE_FILL,
        BLUE,
        1.35,
        name="Module 1 card",
    )
    add_badge(slide, 2.42, 1.08, 1, BLUE, "Module 1 badge")
    add_text(
        slide,
        2.73,
        1.35,
        1.74,
        0.66,
        "Dialogue Feature\nAnalysis",
        size=13.1,
        bold=True,
        font=FONT_DISPLAY,
        name="Module 1 title",
    )
    add_text(
        slide,
        2.73,
        2.17,
        1.74,
        0.30,
        "states  •  strategies  •  transitions",
        size=7.7,
        name="Module 1 details",
    )

    add_box(
        slide,
        5.18,
        1.16,
        2.22,
        1.58,
        BLUE_FILL,
        BLUE,
        1.35,
        name="Module 2 card",
    )
    add_badge(slide, 5.08, 1.08, 2, BLUE, "Module 2 badge")
    add_text(
        slide,
        5.38,
        1.35,
        1.82,
        0.66,
        "Bayesian Dialogue\nModeling",
        size=13.1,
        bold=True,
        font=FONT_DISPLAY,
        name="Module 2 title",
    )
    add_text(
        slide,
        5.38,
        2.14,
        1.82,
        0.38,
        "transition & observation\nprobabilities",
        size=7.9,
        name="Module 2 details",
    )

    add_box(slide, 0.55, 4.31, 1.92, 2.15, WHITE, LINE, 1.35, name="Large corpus")
    add_layers_icon(slide, 1.34, 4.51, MUTED, "Candidate pool icon")
    add_text(
        slide,
        0.79,
        4.93,
        1.44,
        0.26,
        "CANDIDATE POOL",
        size=8.1,
        name="Large corpus role",
    )
    add_text(
        slide,
        0.75,
        5.22,
        1.52,
        0.95,
        "Large, general-\npurpose dialogue\ncorpus",
        size=13.1,
        bold=True,
        font=FONT_DISPLAY,
        name="Large corpus label",
    )

    add_box(
        slide,
        3.30,
        4.55,
        2.20,
        1.55,
        TEAL_FILL,
        TEAL,
        1.35,
        name="Module 3 card",
    )
    add_badge(slide, 3.20, 4.47, 3, TEAL, "Module 3 badge")
    add_text(
        slide,
        3.53,
        4.72,
        1.74,
        0.68,
        "Style Compatibility\nScoring",
        size=13.1,
        bold=True,
        font=FONT_DISPLAY,
        name="Module 3 title",
    )
    add_text(
        slide,
        3.56,
        5.59,
        1.68,
        0.27,
        "Score(rₜ) ∈ [0, 1]",
        size=9.2,
        font=FONT_MATH,
        name="Module 3 details",
    )

    add_box(
        slide,
        6.18,
        4.55,
        3.04,
        1.55,
        ORANGE_FILL,
        ORANGE,
        1.35,
        name="Module 4 card",
    )
    add_badge(slide, 6.08, 4.47, 4, ORANGE, "Module 4 badge")
    add_text(
        slide,
        6.43,
        4.70,
        2.54,
        0.72,
        "Preference Data Selection\n& DPO Adaptation",
        size=13.0,
        bold=True,
        font=FONT_DISPLAY,
        name="Module 4 title",
    )
    add_text(
        slide,
        6.74,
        5.58,
        1.92,
        0.27,
        "chosen  ≻  rejected",
        size=9.1,
        font=FONT_MATH,
        name="Module 4 details",
    )

    add_box(
        slide,
        10.23,
        4.67,
        2.22,
        1.32,
        INK,
        INK,
        1.2,
        name="Output LLM card",
    )
    chip = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.FLOWCHART_MAGNETIC_DISK,
        Inches(11.15),
        Inches(4.85),
        Inches(0.34),
        Inches(0.30),
    )
    set_shape_name(chip, "Output LLM chip icon")
    chip.fill.solid()
    chip.fill.fore_color.rgb = rgb(WHITE)
    chip.line.color.rgb = rgb(WHITE)
    add_text(
        slide,
        10.47,
        5.20,
        1.74,
        0.57,
        "Style-adapted\nlocal LLM",
        size=13.0,
        color=WHITE,
        bold=True,
        font=FONT_DISPLAY,
        name="Output LLM label",
    )

    # Main flows. Labels sit in dedicated whitespace above each connector.
    add_connector(slide, 2.17, 1.99, 2.52, 1.99, width=1.6, name="Small corpus to Module 1")
    add_connector(slide, 4.68, 1.99, 5.18, 1.99, width=1.6, name="Module 1 to Module 2")
    add_line_label(slide, 4.66, 1.66, 0.58, 0.25, "features", name="Dialogue features label")

    add_polyline(
        slide,
        [(6.29, 2.74), (6.29, 3.31), (4.40, 3.31), (4.40, 4.55)],
        color=BLUE,
        width=1.6,
        name="Bayesian model route",
    )
    add_line_label(
        slide,
        4.55,
        3.12,
        1.55,
        0.26,
        "Bayesian dialogue model",
        color=BLUE,
        name="Bayesian model label",
    )

    add_connector(slide, 2.47, 5.32, 3.30, 5.32, color=MUTED, width=1.5, name="Candidates to scoring")
    add_line_label(
        slide,
        2.34,
        5.01,
        1.12,
        0.24,
        "candidate responses",
        size=7.4,
        name="Candidate responses label",
    )
    add_connector(slide, 5.50, 5.32, 6.18, 5.32, color=TEAL, width=1.6, name="Scoring to preference")
    add_line_label(
        slide,
        5.32,
        5.01,
        1.04,
        0.24,
        "scored candidates",
        color=TEAL,
        size=7.4,
        name="Scored candidates label",
    )
    add_connector(slide, 9.22, 5.32, 10.23, 5.32, color=ORANGE, width=1.6, name="DPO to output LLM")
    add_line_label(
        slide,
        9.39,
        5.01,
        0.78,
        0.24,
        "LoRA–DPO",
        color=ORANGE,
        size=7.8,
        name="LoRA DPO label",
    )
    return slide


def add_state_circle(slide, x, y, label, name):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.38), Inches(0.38)
    )
    set_shape_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb(BLUE)
    shape.line.width = Pt(1)
    add_text(
        slide,
        x,
        y,
        0.38,
        0.38,
        label,
        size=8.2,
        font=FONT_MATH,
        name=f"{name} label",
    )


def build_scoring_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)

    panel_y = 1.22
    panel_h = 5.88
    prior_x, predict_x, observe_x, update_x, score_x = 0.26, 2.68, 5.10, 7.52, 10.82
    small_w, update_w = 2.20, 3.08

    add_panel(slide, prior_x, panel_y, small_w, panel_h, 1, "PRIOR STATE", BLUE, "Prior")
    add_panel(slide, predict_x, panel_y, small_w, panel_h, 2, "PREDICT", BLUE, "Predict")
    add_panel(slide, observe_x, panel_y, small_w, panel_h, 3, "LABEL RESPONSE", ORANGE, "Observe")
    add_panel(slide, update_x, panel_y, update_w, panel_h, 4, "BAYESIAN UPDATE", TEAL, "Update")
    add_panel(slide, score_x, panel_y, small_w, panel_h, 5, "STYLE SCORE", TEAL, "Score")

    # Prior state distribution
    add_text(
        slide,
        prior_x + 0.25,
        1.71,
        1.70,
        0.29,
        "state belief bₜ₋₁(s)",
        size=8.9,
        color=MUTED,
        font=FONT_MATH,
        name="Prior distribution label",
    )
    add_connector(slide, 1.20, 2.20, 1.20, 6.45, color=LINE, width=1.0, arrow=False, name="Prior chart axis")
    prior_rows = [
        ("Emotion\ndisclosure", 2.43, 0.48, "0.30", "Prior disclosure"),
        ("Emotion\nprocessing", 3.65, 0.72, "0.45", "Prior processing"),
        ("Considering\nsolutions", 4.87, 0.40, "0.25", "Prior solutions"),
    ]
    for label, y, width, value, name in prior_rows:
        add_text(
            slide,
            prior_x + 0.07,
            y - 0.20,
            0.82,
            0.58,
            label,
            size=8.4,
            align=PP_ALIGN.RIGHT,
            name=f"{name} label",
        )
        add_bar(slide, 1.25, y, width, BLUE, value, name)

    # Prediction
    add_chip(
        slide,
        predict_x + 0.18,
        1.77,
        1.84,
        0.34,
        "transition model  P(sₜ | sₜ₋₁)",
        BLUE,
        BLUE_FILL,
        size=8.5,
        name="Transition model",
    )
    add_state_circle(slide, 3.12, 2.45, "s₁", "State 1")
    add_state_circle(slide, 3.58, 2.45, "s₂", "State 2")
    add_state_circle(slide, 4.04, 2.45, "s₃", "State 3")
    add_connector(slide, 3.50, 2.64, 3.58, 2.64, color=BLUE, width=1.0, name="State 1 to State 2")
    add_connector(slide, 3.96, 2.64, 4.04, 2.64, color=BLUE, width=1.0, name="State 2 to State 3")
    add_text(
        slide,
        predict_x + 0.18,
        3.02,
        1.84,
        1.20,
        "b̂ₜ(sₜ) = ∑  P(sₜ | sₜ₋₁)\n            sₜ₋₁∈S\n            · bₜ₋₁(sₜ₋₁)",
        size=11.1,
        font=FONT_MATH,
        name="Prediction equation",
    )
    add_chip(
        slide,
        predict_x + 0.17,
        5.47,
        1.86,
        0.43,
        "predicted distribution  b̂ₜ(sₜ)",
        BLUE,
        BLUE_FILL,
        size=8.5,
        name="Predicted distribution",
    )

    # Response labeling
    add_box(
        slide,
        observe_x + 0.25,
        1.88,
        1.70,
        0.72,
        WHITE,
        ORANGE,
        1.05,
        name="Candidate response box",
    )
    add_text(
        slide,
        observe_x + 0.37,
        2.00,
        1.46,
        0.48,
        "“That sounds really\ndifficult.”",
        size=10.0,
        font=FONT,
        name="Candidate response",
    )
    add_text(
        slide,
        observe_x + 0.32,
        2.68,
        1.56,
        0.25,
        "candidate response  rₜ",
        size=8.4,
        color=MUTED,
        font=FONT_MATH,
        name="Candidate response variable",
    )
    llm = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(observe_x + 0.80),
        Inches(3.11),
        Inches(0.60),
        Inches(0.60),
    )
    set_shape_name(llm, "LLM classifier")
    llm.fill.solid()
    llm.fill.fore_color.rgb = rgb(ORANGE_FILL)
    llm.line.color.rgb = rgb(ORANGE)
    llm.line.width = Pt(1.1)
    add_text(
        slide,
        observe_x + 0.80,
        3.11,
        0.60,
        0.60,
        "LLM",
        size=10.2,
        color=ORANGE,
        bold=True,
        name="LLM classifier label",
    )
    add_connector(slide, observe_x + 1.10, 2.60, observe_x + 1.10, 3.11, color=ORANGE, width=1.2, name="Response to LLM")
    add_chip(
        slide,
        observe_x + 0.23,
        5.39,
        1.74,
        0.61,
        "observed label  oₜ\nEmpathy / validation",
        ORANGE,
        ORANGE_FILL,
        size=8.7,
        name="Observed label",
    )
    add_connector(slide, observe_x + 1.10, 3.71, observe_x + 1.10, 5.39, color=ORANGE, width=1.2, name="LLM to observed label")

    # Bayesian update
    add_chip(
        slide,
        update_x + 0.26,
        1.79,
        1.18,
        0.36,
        "observation  oₜ",
        ORANGE,
        ORANGE_FILL,
        size=8.6,
        name="Observation input port",
    )
    add_chip(
        slide,
        update_x + 1.69,
        1.79,
        1.12,
        0.36,
        "prediction  b̂ₜ",
        BLUE,
        BLUE_FILL,
        size=8.6,
        name="Prediction input port",
    )
    add_text(
        slide,
        update_x + 0.20,
        2.39,
        2.68,
        1.13,
        "bₜ(sₜ) =  P(oₜ | sₜ) × b̂ₜ(sₜ)\n               ───────────────\n               ∑ₛ′∈S P(oₜ | s′) b̂ₜ(s′)",
        size=13.0,
        font=FONT_MATH,
        name="Bayesian update equation",
    )
    add_connector(slide, update_x + 0.85, 2.15, update_x + 1.10, 2.47, color=ORANGE, width=1.2, name="Observation to likelihood term")
    add_connector(slide, update_x + 2.25, 2.15, update_x + 2.17, 2.47, color=BLUE, width=1.2, name="Prediction to prior term")
    add_text(
        slide,
        update_x + 0.38,
        3.81,
        2.32,
        0.28,
        "posterior state belief  bₜ(s)",
        size=8.7,
        color=MUTED,
        font=FONT_MATH,
        name="Posterior distribution label",
    )
    add_connector(slide, update_x + 1.55, 4.26, update_x + 1.55, 6.41, color=LINE, width=1.0, arrow=False, name="Posterior chart axis")
    posterior_rows = [
        ("Disclosure", 4.53, 0.19, "0.12", MUTED, "Posterior disclosure"),
        ("Processing", 5.18, 1.03, "0.66", TEAL, "Posterior processing"),
        ("Solutions", 5.83, 0.34, "0.22", TEAL, "Posterior solutions"),
    ]
    for label, y, width, value, color, name in posterior_rows:
        add_text(
            slide,
            update_x + 0.28,
            y - 0.06,
            1.14,
            0.28,
            label,
            size=8.5,
            color=color if color == TEAL else INK,
            align=PP_ALIGN.RIGHT,
            name=f"{name} label",
        )
        add_bar(slide, update_x + 1.61, y, width, color, value, name)

    # Score
    add_text(
        slide,
        score_x + 0.18,
        1.82,
        1.84,
        0.78,
        "Score(rₜ) =  ∑  bₜ(s)\n                  s∈Sₜₐᵣɡₑₜ",
        size=11.4,
        font=FONT_MATH,
        name="Style score equation",
    )
    add_chip(
        slide,
        score_x + 0.25,
        2.87,
        1.70,
        0.69,
        "Sₜₐᵣɡₑₜ\nProcessing + Solutions",
        TEAL,
        TEAL_FILL,
        size=8.9,
        name="Target states",
    )
    add_text(
        slide,
        score_x + 0.40,
        3.88,
        1.40,
        0.33,
        "0.66 + 0.22",
        size=12.0,
        font=FONT_MATH,
        name="Score arithmetic",
    )
    add_text(
        slide,
        score_x + 0.47,
        4.53,
        1.26,
        0.58,
        "0.88",
        size=27,
        color=TEAL,
        bold=True,
        font=FONT_MATH,
        name="Style compatibility score",
    )
    add_text(
        slide,
        score_x + 0.36,
        5.22,
        1.48,
        0.28,
        "style compatibility",
        size=8.6,
        color=MUTED,
        name="Style compatibility label",
    )

    # Non-overlapping variable rails above the five panels
    add_polyline(
        slide,
        [(4.71, 5.69), (4.96, 5.69), (4.96, 0.28), (9.77, 0.28), (9.77, 1.79)],
        color=BLUE,
        width=1.45,
        name="Predicted distribution rail",
    )
    add_line_label(
        slide,
        6.85,
        0.13,
        1.86,
        0.28,
        "predicted state distribution",
        color=BLUE,
        size=8.3,
        name="Predicted distribution rail label",
    )
    add_polyline(
        slide,
        [(7.07, 5.70), (7.36, 5.70), (7.36, 0.72), (8.37, 0.72), (8.37, 1.79)],
        color=ORANGE,
        width=1.45,
        name="Observation rail",
    )
    add_line_label(
        slide,
        6.92,
        0.58,
        1.66,
        0.28,
        "response-derived observation",
        color=ORANGE,
        size=8.1,
        name="Observation rail label",
    )

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    # Remove the default first slide only if a template unexpectedly contains one.
    while prs.slides:
        slide_id = prs.slides._sldIdLst[0]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        del prs.slides._sldIdLst[0]
    build_architecture_slide(prs)
    build_scoring_slide(prs)
    prs.core_properties.title = "BASiS Figures 1 and 2 — Editable PowerPoint"
    prs.core_properties.subject = "Editable architecture and Bayesian scoring diagrams"
    prs.core_properties.author = "BASiS paper authors"
    prs.core_properties.comments = (
        "All visible elements are native PowerPoint shapes, text boxes, and connectors."
    )
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
